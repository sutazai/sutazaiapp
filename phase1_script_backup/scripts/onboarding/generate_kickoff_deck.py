#!/usr/bin/env python3
"""
Generate onboarding deck `docs/onboarding/kickoff_deck_v1.pptx` from verified docs.

Requirements:
  pip install python-pptx

Usage:
  python scripts/onboarding/generate_kickoff_deck.py

Notes:
  - Reads `docs/onboarding/kickoff_overview.md` and creates concise slides.
  - No external network access required.
"""
from __future__ import annotations

import sys
from pathlib import Path


def main() -> int:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:  # pragma: no cover
        print("python-pptx is required: pip install python-pptx", file=sys.stderr)
        print(f"Import error: {e}", file=sys.stderr)
        return 2

    src = Path("docs/onboarding/kickoff_overview.md")
    dst = Path("docs/onboarding/kickoff_deck_v1.pptx")
    text = src.read_text(encoding="utf-8") if src.exists() else "Perfect Jarvis Onboarding"

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Perfect Jarvis — Kickoff Overview"
    slide.placeholders[1].text = "Architecture, Ownership, and Health Verification"

    # Technology Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Technology Stack"
    body = slide.shapes.placeholders[1].text_frame
    for line in text.splitlines():
        if line.startswith("- ") and ("Backend:" in line or "Observability:" in line or "Data:" in line or "Vector DB" in line or "Service discovery:" in line or "Frontend:" in line):
            p = body.add_paragraph()
            p.text = line[2:]

    # Ownership
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Ownership"
    body = slide.shapes.placeholders[1].text_frame
    for line in text.splitlines():
        if line.startswith("- ") and ("DevOps" in line or "Observability" in line or "Backend" in line or "Frontend" in line or "Security" in line):
            p = body.add_paragraph(); p.text = line[2:]

    # Agenda
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Onboarding Agenda"
    body = slide.shapes.placeholders[1].text_frame
    agenda = False
    for line in text.splitlines():
        if line.strip().startswith("## Onboarding Meeting Agenda"):
            agenda = True
            continue
        if agenda:
            if line.strip().startswith("## "):
                break
            if line.strip().startswith("- "):
                p = body.add_paragraph(); p.text = line.strip()[2:]

    prs.save(dst)
    print(f"Wrote {dst}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

